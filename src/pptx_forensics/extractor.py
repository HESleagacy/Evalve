"""The authoritative OOXML package parser.

This module deliberately does not use ``python-pptx`` to discover package
parts. The ZIP and relationship layers are parsed directly so unsupported
PresentationML is still retained and visible in the report.
"""

from __future__ import annotations

from collections import defaultdict
from dataclasses import asdict
import hashlib
import importlib.metadata
import posixpath
from pathlib import Path, PurePosixPath
import re
import shutil
from typing import Any, Iterable
import zipfile

from defusedxml import ElementTree as SafeET

from .geometry import IDENTITY, child_transform_matrix, geometry_for_object
from .diagrams import add_native_diagram_evidence
from .models import (
    DECKIR_SCHEMA,
    DECKIR_SCHEMA_VERSION,
    DeckIR,
    ExtractionReport,
    MediaRecord,
    PartRecord,
    RelationshipRecord,
    SlideRecord,
)
from .semantics import native_semantics, resolve_style
from .visual import add_native_visual_evidence

VERSION = "0.5.0"
CONTENT_TYPES = "[Content_Types].xml"
RELATIONSHIP_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
_SLIDE_NUMBER = re.compile(r"^ppt/slides/slide(\d+)\.xml$")


class ExtractionError(ValueError):
    """Raised when the input is not a safe, readable PPTX package."""


def _sha256(data: bytes) -> str:
    return hashlib.sha256(data).hexdigest()


def _local_name(tag: str) -> str:
    return tag.rsplit("}", 1)[-1]


def _attr(element: Any, local_name: str) -> str | None:
    for key, value in element.attrib.items():
        if _local_name(key) == local_name:
            return value
    return None


def _parse_xml(data: bytes, part: str) -> Any:
    try:
        return SafeET.fromstring(data)
    except Exception as exc:
        raise ExtractionError(f"Invalid XML in package part {part}: {exc}") from exc


def _safe_part_name(name: str) -> str:
    if not name or name.startswith("/") or "\\" in name:
        raise ExtractionError(f"Unsafe package part name: {name!r}")
    path = PurePosixPath(name)
    if any(piece in ("", ".", "..") for piece in path.parts):
        raise ExtractionError(f"Unsafe package part name: {name!r}")
    return name


def _relationship_source(name: str) -> str:
    path = PurePosixPath(name)
    if path.name == ".rels" and path.parent.name == "_rels":
        owner = path.parent.parent
        return "" if str(owner) == "." else str(owner)
    if path.parent.name != "_rels":
        raise ExtractionError(f"Invalid relationships part name: {name}")
    return str(path.parent.parent / path.stem)


def _resolve_target(source: str, target: str, mode: str | None) -> str | None:
    if (mode and mode.lower() == "external") or not target:
        return None
    if target.startswith("/"):
        resolved = posixpath.normpath(target.lstrip("/"))
    else:
        base = PurePosixPath(source).parent if source else PurePosixPath(".")
        resolved = posixpath.normpath(posixpath.join(str(base), target))
    if resolved == "." or resolved.startswith("../") or "/../" in resolved:
        return None
    return resolved.lstrip("./")


def _dependency_versions() -> dict[str, str]:
    versions: dict[str, str] = {}
    for distribution in ("defusedxml", "lxml", "python-pptx"):
        try:
            versions[distribution] = importlib.metadata.version(distribution)
        except importlib.metadata.PackageNotFoundError:
            versions[distribution] = "not-installed"
    return versions


def _content_types(raw: bytes | None) -> dict[str, str]:
    if raw is None:
        return {}
    root = _parse_xml(raw, CONTENT_TYPES)
    result: dict[str, str] = {}
    for item in root:
        name = _local_name(item.tag)
        if name == "Override" and item.get("PartName") and item.get("ContentType"):
            result[item.get("PartName").lstrip("/")] = item.get("ContentType")
        elif name == "Default" and item.get("Extension") and item.get("ContentType"):
            result[f"*.{item.get('Extension').lower()}"] = item.get("ContentType")
    return result


def _part_content_type(name: str, content_types: dict[str, str]) -> str:
    if name in content_types:
        return content_types[name]
    suffix = Path(name).suffix.lower().lstrip(".")
    return content_types.get(f"*.{suffix}", "application/octet-stream")


def _relationships(
    names_to_bytes: dict[str, bytes], warnings: list[str]
) -> tuple[list[RelationshipRecord], dict[str, list[RelationshipRecord]]]:
    records: list[RelationshipRecord] = []
    by_source: dict[str, list[RelationshipRecord]] = defaultdict(list)
    for name in sorted(names_to_bytes):
        if not name.endswith(".rels") or "/_rels/" not in f"/{name}":
            continue
        try:
            source = _relationship_source(name)
            root = _parse_xml(names_to_bytes[name], name)
        except ExtractionError as exc:
            warnings.append(str(exc))
            continue
        for relation in root:
            if _local_name(relation.tag) != "Relationship":
                continue
            target = relation.get("Target", "")
            item = RelationshipRecord(
                source_part=source,
                relationship_id=relation.get("Id", ""),
                relationship_type=relation.get("Type", ""),
                target=target,
                target_mode=relation.get("TargetMode"),
                resolved_target=_resolve_target(source, target, relation.get("TargetMode")),
            )
            records.append(item)
            by_source[source].append(item)
    records.sort(key=lambda item: (item.source_part, item.relationship_id, item.relationship_type, item.target))
    return records, by_source


def _rels_to(by_source: dict[str, list[RelationshipRecord]], source: str, suffix: str) -> str | None:
    for relation in by_source.get(source, []):
        if relation.relationship_type.endswith(suffix):
            return relation.resolved_target
    return None


def _texts(root: Any) -> list[str]:
    return [element.text for element in root.iter() if _local_name(element.tag) == "t" and element.text]


def _hyperlinks(root: Any, relations: Iterable[RelationshipRecord]) -> list[dict[str, Any]]:
    relation_map = {item.relationship_id: item for item in relations}
    found: list[dict[str, Any]] = []
    for element in root.iter():
        kind = _local_name(element.tag)
        if kind not in {"hlinkClick", "hlinkHover"}:
            continue
        relationship_id = _attr(element, "id")
        relation = relation_map.get(relationship_id or "")
        found.append(
            {
                "kind": kind,
                "relationship_id": relationship_id,
                "target": relation.target if relation else None,
                "resolved_target": relation.resolved_target if relation else None,
                "action": element.get("action"),
            }
        )
    return found


def _alt_text(root: Any) -> list[dict[str, str]]:
    result: list[dict[str, str]] = []
    for element in root.iter():
        if _local_name(element.tag) != "cNvPr":
            continue
        descr, title = element.get("descr"), element.get("title")
        if descr or title:
            result.append({"id": element.get("id", ""), "name": element.get("name", ""), "descr": descr or "", "title": title or ""})
    return result


def _animations(root: Any) -> list[dict[str, Any]]:
    animation_names = {"timing", "par", "seq", "anim", "set", "animEffect", "animMotion", "cmd"}
    result: list[dict[str, Any]] = []
    for element in root.iter():
        name = _local_name(element.tag)
        if name in animation_names:
            result.append({"element": name, "attributes": dict(element.attrib)})
    return result


def _first_descendant(root: Any, local_name: str) -> Any | None:
    return next((item for item in root.iter() if _local_name(item.tag) == local_name), None)


def _slide_dimensions(root: Any | None) -> tuple[float, float]:
    if root is not None:
        size = _first_descendant(root, "sldSz")
        if size is not None and size.get("cx") and size.get("cy"):
            return float(size.get("cx")), float(size.get("cy"))
    # ISO/IEC 29500 default widescreen-independent slide size (10 x 7.5 in).
    return 9_144_000.0, 6_858_000.0


def _xml_path(parent_path: str, element: Any, siblings: list[Any]) -> str:
    local = _local_name(element.tag)
    same_name = [item for item in siblings if _local_name(item.tag) == local]
    return f"{parent_path}/{local}[{same_name.index(element) + 1}]"


def _shape_id_element(element: Any) -> Any | None:
    return _first_descendant(element, "cNvPr")


def _placeholder_geometry(element: Any, root: Any | None) -> Any | None:
    placeholder = _first_descendant(element, "ph")
    if placeholder is None or root is None:
        return None
    wanted_type = placeholder.get("type", "body")
    wanted_idx = placeholder.get("idx")
    for candidate in root.iter():
        if _local_name(candidate.tag) not in {"sp", "pic", "graphicFrame"}:
            continue
        candidate_placeholder = _first_descendant(candidate, "ph")
        if candidate_placeholder is None or candidate_placeholder.get("type", "body") != wanted_type:
            continue
        if wanted_idx is not None and candidate_placeholder.get("idx") != wanted_idx:
            continue
        if _first_descendant(candidate, "xfrm") is not None:
            return candidate
    return None


def _object_type(element: Any) -> str | None:
    local = _local_name(element.tag)
    if local == "sp":
        return "text" if _first_descendant(element, "txBody") is not None else "shape"
    if local == "pic":
        return "image"
    if local == "cxnSp":
        return "connector"
    if local == "grpSp":
        return "group"
    if local == "graphicFrame":
        graphic_data = _first_descendant(element, "graphicData")
        uri = graphic_data.get("uri", "") if graphic_data is not None else ""
        descendants = {_local_name(item.tag) for item in element.iter()}
        if "table" in descendants or uri.endswith("/table"):
            return "table"
        if "chart" in descendants or "chart" in uri:
            return "chart"
        if "relIds" in descendants or "diagram" in uri:
            return "smartart"
        return "shape"
    if local in {"contentPart", "oleObj"}:
        return "shape"
    return None


def _native_shape_type(element: Any, semantic_type: str | None) -> str | None:
    """Preserve the OOXML/Python-pptx shape subtype beside the semantic type."""
    local = _local_name(element.tag)
    if local == "sp":
        if _first_descendant(element, "ph") is not None:
            return "PLACEHOLDER"
        if _first_descendant(element, "custGeom") is not None:
            return "FREEFORM"
        return "TEXT_BOX" if _first_descendant(element, "txBody") is not None else "AUTO_SHAPE"
    if local == "pic":
        return "PICTURE"
    if local == "cxnSp":
        return "CONNECTOR"
    if local == "grpSp":
        return "GROUP"
    return {
        "table": "TABLE",
        "chart": "CHART",
        "smartart": "SMARTART",
    }.get(semantic_type, "GRAPHIC_FRAME" if local == "graphicFrame" else semantic_type)


def _object_relationship_ids(element: Any) -> list[str]:
    result: list[str] = []
    for descendant in element.iter():
        for key, value in descendant.attrib.items():
            namespace, _, local = key.partition("}")
            if namespace.endswith("relationships") and local in {"id", "embed", "link", "dm", "lo", "qs"}:
                result.append(value)
    return result


def _object_style(element: Any) -> dict[str, Any]:
    placeholder = _first_descendant(element, "ph")
    style: dict[str, Any] = {}
    if placeholder is not None and placeholder.get("type"):
        style["placeholder_type"] = placeholder.get("type")
    if placeholder is not None and placeholder.get("idx"):
        style["placeholder_idx"] = placeholder.get("idx")
    transform = _first_descendant(element, "xfrm")
    if transform is not None and transform.get("rot"):
        try:
            style["rotation_degrees"] = float(transform.get("rot")) / 60_000.0
        except (TypeError, ValueError):
            pass
    return style


def _object_text(element: Any) -> str:
    values = [value.strip() for value in _texts(element) if value.strip()]
    return " ".join(values)


def _native_source(part: str | None, xml_path: str | None, confidence: float = 1.0) -> dict[str, Any]:
    return {
        "layer": "native_ooxml",
        "xml_part": part,
        "xml_path": xml_path,
        "confidence": confidence,
    }


def _slide_objects(
    root: Any,
    slide_id: str,
    part: str,
    slide_relations: Iterable[RelationshipRecord],
    slide_width: float,
    slide_height: float,
    asset_ids: dict[str, str],
    package_parts: dict[str, bytes],
    layout_root: Any | None,
    master_root: Any | None,
    theme_root: Any | None,
) -> list[dict[str, Any]]:
    shape_tree = _first_descendant(root, "spTree")
    if shape_tree is None:
        return []
    relation_ids = {
        item.relationship_id: f"{part}:{item.relationship_id}" for item in slide_relations
    }
    objects: list[dict[str, Any]] = []
    ordinal = 0

    def visit(
        container: Any,
        container_path: str,
        parent_id: str | None,
        parent_matrix: tuple[float, float, float, float, float, float] = IDENTITY,
        transform_chain: list[str] | None = None,
    ) -> None:
        nonlocal ordinal
        ancestor_chain = transform_chain or []
        children = list(container)
        for child in children:
            kind = _object_type(child)
            child_path = _xml_path(container_path, child, children)
            if kind is None:
                continue
            ordinal += 1
            object_id = f"{slide_id}-shape-{ordinal:02d}"
            bbox, geometry = geometry_for_object(
                child,
                slide_width,
                slide_height,
                parent_matrix,
                [*ancestor_chain, object_id],
                _placeholder_geometry(child, layout_root) or _placeholder_geometry(child, master_root),
            )
            shape_id = _shape_id_element(child)
            object_relationships = [
                relation_ids.get(value, f"{part}:{value}")
                for value in _object_relationship_ids(child)
            ]
            object_record: dict[str, Any] = {
                "id": object_id,
                "slide_id": slide_id,
                "parent_id": parent_id,
                "type": kind,
                "shape_type": _native_shape_type(child, kind),
                "bbox": bbox,
                "z_order": ordinal,
                "text": _object_text(child),
                "style": _object_style(child),
                "geometry": geometry,
                "relationships": object_relationships,
                "source": _native_source(part, child_path),
            }
            raw_style, resolved_style, inherited_from = resolve_style(
                child,
                layout_root,
                master_root,
                theme_root,
            )
            object_record["raw_style"] = raw_style
            object_record["resolved_style"] = resolved_style
            object_record["inherited_from"] = inherited_from
            object_record.update(
                native_semantics(
                    child,
                    kind,
                    slide_relations,
                    package_parts,
                    asset_ids,
                )
            )
            if shape_id is not None:
                object_record["native_id"] = shape_id.get("id")
                object_record["name"] = shape_id.get("name", "")
            for relationship_id in _object_relationship_ids(child):
                relation = next(
                    (item for item in slide_relations if item.relationship_id == relationship_id),
                    None,
                )
                if relation and relation.resolved_target in asset_ids:
                    object_record["asset_id"] = asset_ids[relation.resolved_target]
                    if relation.resolved_target.startswith("ppt/embeddings/"):
                        object_record["embedded"] = True
                    break
            objects.append(object_record)
            if kind == "group":
                visit(
                    child,
                    child_path,
                    object_id,
                    child_transform_matrix(child, parent_matrix),
                    [*ancestor_chain, object_id],
                )

    visit(shape_tree, "/sld[1]", None)
    return objects


def _python_pptx_summary(source: Path) -> dict[str, Any]:
    try:
        import pptx  # type: ignore
    except ImportError:
        return {"available": False, "reason": "python-pptx is not installed"}
    try:
        presentation = pptx.Presentation(str(source))
        slides: list[dict[str, Any]] = []
        for number, slide in enumerate(presentation.slides, 1):
            shapes: list[dict[str, Any]] = []
            for shape in slide.shapes:
                item: dict[str, Any] = {"name": shape.name, "shape_type": str(shape.shape_type)}
                if hasattr(shape, "text") and shape.text:
                    item["text"] = shape.text
                shapes.append(item)
            slides.append({"number": number, "shapes": shapes})
        return {"available": True, "version": getattr(pptx, "__version__", "unknown"), "slides": slides}
    except Exception as exc:  # enrichment must not hide package evidence
        return {"available": True, "error": str(exc)}


def _comments(names_to_bytes: dict[str, bytes]) -> list[dict[str, Any]]:
    result: list[dict[str, Any]] = []
    for name, raw in sorted(names_to_bytes.items()):
        if "/comments/" not in f"/{name}" or not name.endswith(".xml"):
            continue
        root = _parse_xml(raw, name)
        for comment in root.iter():
            if _local_name(comment.tag) not in {"comment", "cm"}:
                continue
            result.append({"part": name, "author_id": _attr(comment, "authorId"), "text": "".join(_texts(comment))})
    return result


def extract_pptx(
    source: str | Path,
    evidence_dir: str | Path | None = None,
    *,
    include_visual_evidence: bool = True,
    include_native_diagrams: bool = True,
) -> ExtractionReport:
    """Extract a PPTX without discarding unsupported OOXML.

    When ``evidence_dir`` is supplied, the original archive is copied as
    ``original.pptx`` and every ZIP member is retained below ``parts/``.
    """
    source_path = Path(source).expanduser().resolve()
    if not source_path.is_file():
        raise ExtractionError(f"PPTX file does not exist: {source_path}")
    source_bytes = source_path.read_bytes()
    warnings: list[str] = []
    try:
        archive = zipfile.ZipFile(source_path)
    except (zipfile.BadZipFile, OSError) as exc:
        raise ExtractionError(f"Not a readable PPTX ZIP package: {exc}") from exc

    with archive:
        infos = archive.infolist()
        names_to_bytes: dict[str, bytes] = {}
        for info in infos:
            name = _safe_part_name(info.filename)
            if name in names_to_bytes:
                warnings.append(f"Duplicate ZIP member retained once: {name}")
                continue
            if info.is_dir():
                continue
            try:
                names_to_bytes[name] = archive.read(info)
            except (RuntimeError, OSError, zipfile.BadZipFile) as exc:
                raise ExtractionError(f"Could not read package part {name}: {exc}") from exc

    if CONTENT_TYPES not in names_to_bytes:
        warnings.append("Missing [Content_Types].xml")
    content_types = _content_types(names_to_bytes.get(CONTENT_TYPES))
    relationships, by_source = _relationships(names_to_bytes, warnings)
    part_records = [
        PartRecord(
            name=name,
            content_type=_part_content_type(name, content_types),
            size=len(raw),
            compressed_size=next((item.compress_size for item in infos if item.filename == name), len(raw)),
            crc32=f"{next((item.CRC for item in infos if item.filename == name), 0):08x}",
            sha256=_sha256(raw),
            is_xml=name.endswith(".xml") or name.endswith(".rels"),
        )
        for name, raw in sorted(names_to_bytes.items())
    ]
    media = [
        MediaRecord(part=item.name, content_type=item.content_type, size=item.size, sha256=item.sha256)
        for item in part_records
        if item.name.startswith("ppt/media/")
    ]
    asset_parts = [
        item
        for item in part_records
        if item.name.startswith("ppt/media/") or item.name.startswith("ppt/embeddings/")
    ]
    asset_ids = {item.name: f"asset-{index:04d}" for index, item in enumerate(asset_parts, 1)}

    slides_from_relationships: list[str] = []
    presentation = names_to_bytes.get("ppt/presentation.xml")
    presentation_root: Any | None = None
    if presentation:
        presentation_root = _parse_xml(presentation, "ppt/presentation.xml")
        presentation_relations = {item.relationship_id: item for item in by_source.get("ppt/presentation.xml", [])}
        for element in presentation_root.iter():
            if _local_name(element.tag) != "sldId":
                continue
            relationship_id = element.get(f"{{{RELATIONSHIP_NS}}}id")
            relation = presentation_relations.get(relationship_id or "")
            if relation and relation.resolved_target:
                slides_from_relationships.append(relation.resolved_target)
    if not slides_from_relationships:
        slides_from_relationships = [name for name in sorted(names_to_bytes) if _SLIDE_NUMBER.match(name)]

    slides: list[SlideRecord] = []
    canonical_slides: list[dict[str, Any]] = []
    canonical_objects: list[dict[str, Any]] = []
    slide_width, slide_height = _slide_dimensions(presentation_root)
    for fallback_number, part in enumerate(slides_from_relationships, 1):
        raw = names_to_bytes.get(part)
        if raw is None:
            warnings.append(f"Slide relationship points to missing part: {part}")
            continue
        root = _parse_xml(raw, part)
        layout = _rels_to(by_source, part, "/slideLayout")
        master = _rels_to(by_source, layout, "/slideMaster") if layout else None
        theme = _rels_to(by_source, master, "/theme") if master else None
        layout_root = _parse_xml(names_to_bytes[layout], layout) if layout in names_to_bytes else None
        master_root = _parse_xml(names_to_bytes[master], master) if master in names_to_bytes else None
        theme_root = _parse_xml(names_to_bytes[theme], theme) if theme in names_to_bytes else None
        notes_part = _rels_to(by_source, part, "/notesSlide")
        notes = _texts(_parse_xml(names_to_bytes[notes_part], notes_part)) if notes_part in names_to_bytes else []
        slide_text = _texts(root)
        slide_hyperlinks = _hyperlinks(root, by_source.get(part, []))
        slide_alt_text = _alt_text(root)
        slide_animations = _animations(root)
        slides.append(
            SlideRecord(
                number=fallback_number,
                part=part,
                layout_part=layout,
                master_part=master,
                theme_part=theme,
                text=slide_text,
                notes=notes,
                hyperlinks=slide_hyperlinks,
                alt_text=slide_alt_text,
                animations=slide_animations,
            )
        )
        slide_id = f"slide-{fallback_number:02d}"
        slide_objects = _slide_objects(
            root,
            slide_id,
            part,
            by_source.get(part, []),
            slide_width,
            slide_height,
            asset_ids,
            names_to_bytes,
            layout_root,
            master_root,
            theme_root,
        )
        canonical_objects.extend(slide_objects)
        canonical_slides.append(
            {
                "id": slide_id,
                "number": fallback_number,
                "part": part,
                "layout_part": layout,
                "master_part": master,
                "theme_part": theme,
                "text": slide_text,
                "notes": notes,
                "hyperlinks": slide_hyperlinks,
                "alt_text": slide_alt_text,
                "animations": slide_animations,
                "object_ids": [item["id"] for item in slide_objects],
                "slide_reading_order": "unknown",
                "diagram_flow_direction": "unknown",
                "flow_present": None,
                "flow_presence_basis": "undetermined",
                "visual_region_ids": [],
                "visual_evidence_visibility": {
                    "native": "verified",
                    "rendered": "not_requested",
                    "ocr": "not_requested",
                    "vision": "not_requested",
                },
                "source": _native_source(part, "/sld[1]"),
            }
        )

    canonical_relationships = []
    for relation in relationships:
        relation_id = f"{relation.source_part or 'package'}:{relation.relationship_id}"
        canonical_relationships.append(
            {
                "id": relation_id,
                "source_part": relation.source_part,
                "relationship_id": relation.relationship_id,
                "relationship_type": relation.relationship_type,
                "target": relation.target,
                "target_mode": relation.target_mode,
                "resolved_target": relation.resolved_target,
                "source": _native_source(relation.source_part or None, None),
            }
        )
    canonical_assets = [
        {
            "id": asset_ids[item.name],
            "type": "embedded" if item.name.startswith("ppt/embeddings/") else "media",
            "part": item.name,
            "content_type": item.content_type,
            "size": item.size,
            "sha256": item.sha256,
            "source": {
                "layer": "native_ooxml",
                "package_part": item.name,
                "xml_part": None,
                "xml_path": None,
                "confidence": 1.0,
            },
        }
        for item in asset_parts
    ]
    canonical = DeckIR(
        deck={
            "id": "deck",
            "schema": DECKIR_SCHEMA,
            "schema_version": DECKIR_SCHEMA_VERSION,
            "name": source_path.name,
            "source": source_path.name,
            "source_sha256": _sha256(source_bytes),
            "slide_count": len(canonical_slides),
            "slide_size_emu": [slide_width, slide_height],
            "slide_aspect_ratio": round(slide_width / slide_height, 12) if slide_height else None,
            "package_part_count": len(part_records),
        },
        slides=canonical_slides,
        objects=canonical_objects,
        assets=canonical_assets,
        relationships=canonical_relationships,
        rendered_evidence=[],
        ocr_evidence=[],
        vision_evidence=[],
        warnings=warnings,
        provenance={
            "schema": DECKIR_SCHEMA,
            "schema_version": DECKIR_SCHEMA_VERSION,
            "parser_version": VERSION,
            "dependencies": _dependency_versions(),
            "source_name": source_path.name,
            "package_parts": [asdict(item) for item in part_records],
            "authority": {
                "native_ooxml": "authoritative",
                "rendered_cv": "derived",
                "ocr": "derived",
                "vision_model": "probabilistic",
            },
        },
    )
    if include_visual_evidence:
        add_native_visual_evidence(canonical, names_to_bytes)
    if include_native_diagrams:
        add_native_diagram_evidence(canonical)

    report = ExtractionReport(
        source=str(source_path),
        source_sha256=_sha256(source_bytes),
        parser_version=VERSION,
        dependencies=_dependency_versions(),
        package_parts=part_records,
        relationships=relationships,
        slides=slides,
        media=media,
        comments=_comments(names_to_bytes),
        convenience=_python_pptx_summary(source_path),
        warnings=warnings,
        canonical=canonical,
    )
    if evidence_dir is not None:
        destination = Path(evidence_dir).expanduser().resolve()
        destination.mkdir(parents=True, exist_ok=True)
        (destination / "parts").mkdir(exist_ok=True)
        shutil.copyfile(source_path, destination / "original.pptx")
        for name, raw in names_to_bytes.items():
            target = destination / "parts" / name
            target.parent.mkdir(parents=True, exist_ok=True)
            target.write_bytes(raw)
        report.evidence_dir = str(destination)
    return report
